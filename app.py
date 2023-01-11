import streamlit as st
import yfinance as yf
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
from datetime import date
from PIL import Image
import requests
import os
from io import BytesIO



def resize_image(url):
    """function to resize logos while keeping aspect ratio. Accepts URL as an argument and return an image object"""

    # Open the image file
    image = Image.open(requests.get(url, stream=True).raw)

    # if a logo is too high or too wide then make the background container twice as big
    if image.height > 140:
        container_width = 220 * 2
        container_height = 140 * 2

    elif image.width > 220:
        container_width = 220 * 2
        container_height = 140 * 2
    else:
        container_width = 220
        container_height = 140

    # Create a new image with the same aspect ratio as the original image
    new_image = Image.new('RGBA', (container_width, container_height))

    # Calculate the position to paste the image so that it is centered
    x = (container_width - image.width) // 2
    y = (container_height - image.height) // 2

    # Paste the image onto the new image
    new_image.paste(image, (x, y))

    return new_image


def add_image(slide, image, left, top, width):
    """function to add an image to the PowerPoint slide and specify its position and width"""
    slide.shapes.add_picture(image, left=left, top=top, width=width)


# function to replace text in pptx first slide with selected filters
def replace_text(replacements, shapes):
    """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text



def get_stock(ticker, period, interval):
    hist = ticker.history(period=period, interval=interval)
    stock_df = pd.DataFrame(hist).reset_index()

    return stock_df


def plot_graph(df, x, y, title):
    fig = px.line(df, x=x, y=y, template='simple_white',
                        title='<b>{} {}</b>'.format(name, title))
    fig.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

    return fig



def get_financials(df, col_name, metric_name):
    metric = df.loc[col_name]
    metric_df = pd.DataFrame(metric).reset_index()
    metric_df.columns = ['Year', metric_name]

    return metric_df



#set page layout
st.set_page_config(page_icon="🚀", page_title="PowerPoint Generator")

path = os.path.dirname(__file__)
today = date.today()


st.title('Welcome to the PowerPoint Generator!')

st.subheader('This web app uses data from Yahoo Finance to create PowerPoint slides')

user_input = st.text_input(label='Enter company ticker. For example: AAPL for Apple or TSLA for Tesla')

options = ['Stock Price', 'Revenue', 'Cashflow', 'EBITDA']
selected_options = st.multiselect(label='Select metric(s)', options=options)

submit = st.button(label='Generate PowerPoint slides')

# trim user input string
user_input = str(user_input).strip()

if submit and user_input == "":
    st.warning("Please enter company ticker!")

elif submit and user_input != "":
    with st.spinner('Generating awesome slides for you...⏳'):
        try:

            # open presentation template
            pptx = path + '//' + 'template.pptx'
            prs = Presentation(pptx)

            # declare positional variables
            width = Inches(8)
            left = Inches(2.5)
            top = Inches(1)

            ticker = yf.Ticker(user_input)

            # get stock info
            name = ticker.info['shortName']
            logo_url = ticker.info['logo_url']
            sector = ticker.info['sector']
            industry = ticker.info['industry']
            employees = ticker.info['fullTimeEmployees']
            country = ticker.info['country']
            city = ticker.info['city']
            website = ticker.info['website']
            summary = ticker.info['longBusinessSummary']

            # declare pptx variables
            first_slide = prs.slides[0]
            second_slide = prs.slides[1]
            shapes_1 = []
            shapes_2 = []
            index_to_drop = []

            # create lists with shape objects
            for shape in first_slide.shapes:
                shapes_1.append(shape)

            for shape in second_slide.shapes:
                shapes_2.append(shape)

            # initiate a dictionary of placeholders and values to replace
            replaces_1 = {
                '{company}': name,
                '{date}': today}

            replaces_2 = {
                '{c}': name,
                '{s}': sector,
                '{i}': industry,
                '{co}': country,
                '{ci}': city,
                '{ee}': employees,
                '{w}': website,
                '{summary}': summary
            }

            # run the function to replace placeholders with values
            replace_text(replaces_1, shapes_1)
            replace_text(replaces_2, shapes_2)

            # create logo image object
            logo = resize_image(logo_url)
            logo.save('logo.png')
            logo_im = path + '//' + 'logo.png'

            # add plots
            add_image(prs.slides[1], image=logo_im, left=Inches(1.2), width=Inches(2), top=Inches(0.5))
            os.remove('logo.png')


            if len(selected_options) > 0:
                income_df = pd.DataFrame(ticker.income_stmt)
                cashflow_df = pd.DataFrame(ticker.cashflow)

            if 'Stock Price' in selected_options:
                stock_df = get_stock(ticker=ticker, period='5y', interval='1mo')
                stock_fig = plot_graph(df=stock_df, x='Date', y='Open', title='Stock Price USD')

                stock_fig.write_image("stock.png")
                stock_im = path + '//' + 'stock.png'

                add_image(prs.slides[2], image=stock_im, left=left, width=width, top=top)
                os.remove('stock.png')
            else:
                index_to_drop.append(2)

            if "Revenue" in selected_options:
                rev_df = get_financials(df=income_df,col_name='Total Revenue', metric_name='Total Revenue')
                rev_fig = plot_graph(df=rev_df, x='Year', y='Total Revenue', title='Total Revenue USD')

                rev_fig.write_image("rev.png")
                rev_im = path + '//' + 'rev.png'

                add_image(prs.slides[3], image=rev_im, left=left, width=width, top=top)
                os.remove('rev.png')
            else:
                index_to_drop.append(3)

            if "Cashflow" in selected_options:
                cash_df = get_financials(df=cashflow_df, col_name='Operating Cash Flow', metric_name='Operating Cash Flow')
                cashflow_fig = plot_graph(df=cash_df, x='Year', y='Operating Cash Flow', title='Operating Cash Flow USD')

                cashflow_fig.write_image("cashflow.png")
                cashflow_im = path + '//' + 'cashflow.png'

                add_image(prs.slides[4], image=cashflow_im, left=left, width=width, top=top)
                os.remove('cashflow.png')
            else:
                index_to_drop.append(4)

            if "EBITDA" in selected_options:
                ebitda_df = get_financials(df=income_df, col_name='Normalized EBITDA', metric_name='EBITDA')
                ebitda_fig = plot_graph(df=ebitda_df, x='Year', y='EBITDA', title='EBITDA USD')

                ebitda_fig.write_image("ebitda.png")
                ebitda_im = path + '//' + 'ebitda.png'

                add_image(prs.slides[5], image=ebitda_im, left=left, width=width, top=top)
                os.remove('ebitda.png')
            else:
                index_to_drop.append(5)

            if len(selected_options) == 0:
                index_to_drop = [2, 3, 4, 5]
                
            # iterate over index to drop and delete unused slides
            x = 0
            for i in index_to_drop:
                rId = prs.slides._sldIdLst[int(i) - x].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[int(i) - x]
                x += 1

            # create file name
            filename = '{} {}.pptx'.format(name, today)

            # save presentation as binary output
            binary_output = BytesIO()
            prs.save(binary_output)

            # display success message and download button
            st.success('The slides have been generated! :tada:')

            st.download_button(label='Click to download PowerPoint',
                               data=binary_output.getvalue(),
                               file_name=filename)

        except Exception as e:
            st.write(e)
            st.error("Oops, something went wrong! Try again or use a different ticker")


